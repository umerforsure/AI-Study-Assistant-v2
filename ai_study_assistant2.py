# ai_study_assistant2.py
# Final backend with improved T5 summarization and hybrid T5->Qwen refinement (AUTO)
# T5 is used for summaries/notes (project requirement).
# Qwen is used for RAG QA and for refining summaries when needed.

import os
import json
import uuid
import datetime
import re
from pathlib import Path
from typing import List, Dict, Optional

import torch
from transformers import AutoModelForCausalLM, AutoTokenizer, T5ForConditionalGeneration, T5TokenizerFast

# ---------------- CONFIG (change these if you move files) ----------------
LOCAL_QWEN_PATH = r"E:\My Projects\AI Study Assistant\New folder\models--Qwen--Qwen2.5-1.5B-Instruct\snapshots\989aa7980e4cf806f80c7fef2b1adb7bc71aa306"
T5_MODEL_PATH = r"E:\My Projects\AI Study Assistant\New folder\t5_finetuned"
EMBED_MODEL_NAME = "all-MiniLM-L6-v2"
# -----------------------------------------------------------------------

# Optional libs
try:
    from sentence_transformers import SentenceTransformer
except Exception:
    SentenceTransformer = None

try:
    import faiss
except Exception:
    faiss = None

# File extraction optional libs
try:
    from pypdf import PdfReader
except Exception:
    PdfReader = None

try:
    import docx
except Exception:
    docx = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    from PIL import Image
    import pytesseract
except Exception:
    Image = None
    pytesseract = None

# device
device = "cuda" if torch.cuda.is_available() else "cpu"

# --------------------------- Utilities ----------------------------------
def clean_text(s: str) -> str:
    """Normalize whitespace and remove common table separators."""
    if not isinstance(s, str):
        return ""
    s = re.sub(r"[-=]{3,}", " ", s)
    s = re.sub(r"\|+", " ", s)
    s = re.sub(r"\s+", " ", s).strip()
    return s

def chunk_text(text: str, chunk_size: int = 400, overlap: int = 64) -> List[str]:
    text = clean_text(text)
    tokens = text.split()
    if not tokens:
        return []
    chunks = []
    i = 0
    while i < len(tokens):
        chunk = tokens[i:i + chunk_size]
        chunks.append(" ".join(chunk))
        i += chunk_size - overlap
    return chunks

# ---------------------- Embeddings & FAISS helpers ----------------------
_embed = None

def embed_texts(texts: List[str]):
    """Return normalized embeddings (numpy)."""
    global _embed
    if SentenceTransformer is None:
        raise RuntimeError("sentence-transformers not installed")
    if _embed is None:
        _embed = SentenceTransformer(EMBED_MODEL_NAME)
    return _embed.encode(texts, convert_to_numpy=True, normalize_embeddings=True)

def create_faiss_index(dim: int):
    if faiss is None:
        raise RuntimeError("faiss not installed")
    return faiss.IndexFlatIP(dim)

def build_index(chunks: List[str], save_path: str):
    vecs = embed_texts(chunks)
    idx = create_faiss_index(vecs.shape[1])
    idx.add(vecs)
    os.makedirs(os.path.dirname(save_path), exist_ok=True)
    faiss.write_index(idx, save_path)
    return idx

def load_index(path: str):
    if faiss is None:
        raise RuntimeError("faiss not installed")
    if not os.path.exists(path):
        raise FileNotFoundError(f"FAISS index not found at: {path}")
    return faiss.read_index(path)

def load_metas(path: str) -> List[Dict]:
    if not os.path.exists(path):
        raise FileNotFoundError(f"Meta JSON not found at: {path}")
    with open(path, "r", encoding="utf-8") as f:
        return json.load(f)

def search_index(index, query: str, top_k: int = 5):
    qv = embed_texts([query])
    dists, idxs = index.search(qv, top_k)
    return dists[0], idxs[0]

def retrieve_contexts(index, meta_store: List[Dict], question: str, top_k: int = 4, min_score: float = 0.05):
    """
    Return deduped top-k contexts; min_score filters out very-low-similarity results.
    min_score default lowered to 0.05 for your documents.
    """
    dists, idxs = search_index(index, question, top_k=top_k)
    contexts = []
    for dist, i in zip(dists, idxs):
        if i < len(meta_store):
            # dist may be numpy scalar
            try:
                score = float(dist)
            except Exception:
                score = None
            if score is None or score >= min_score:
                contexts.append(meta_store[i]["text"])
    # dedupe by prefix
    seen = set()
    out = []
    for c in contexts:
        key = c[:200].strip()
        if key not in seen:
            seen.add(key)
            out.append(c)
    return out

# --------------------------- File extraction ----------------------------
def extract_text_from_pdf(path: str) -> str:
    if PdfReader is None:
        raise RuntimeError("pypdf not installed")
    reader = PdfReader(path)
    return "\n".join(page.extract_text() or "" for page in reader.pages)

def extract_text_from_docx(path: str) -> str:
    if docx is None:
        raise RuntimeError("python-docx not installed")
    d = docx.Document(path)
    return "\n".join(p.text for p in d.paragraphs if p.text)

def extract_text_from_pptx(path: str) -> str:
    if Presentation is None:
        raise RuntimeError("python-pptx not installed")
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t = shape.text.strip()
                if t:
                    texts.append(t)
    return "\n".join(texts)

def extract_text_from_txt(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()

def extract_text_from_image(path: str) -> str:
    if pytesseract is None or Image is None:
        raise RuntimeError("pytesseract or Pillow not installed")
    img = Image.open(path)
    return pytesseract.image_to_string(img)

# --------------------------- ChatHistory --------------------------------
class ChatHistory:
    def __init__(self, max_items: int = 500):
        self.history = []
        self.max_items = max_items

    def add(self, role: str, text: str):
        self.history.append({
            "role": role,
            "text": text,
            "time": datetime.datetime.utcnow().isoformat()
        })
        if len(self.history) > self.max_items:
            self.history = self.history[-self.max_items:]

    def export_text(self) -> str:
        return "\n".join(f"[{h['time']}] {h['role'].upper()}:\n{h['text']}" for h in self.history)

# --------------------------- T5 (local) ---------------------------------
_t5 = None
_t5_tok = None

def load_t5(model_path: str = T5_MODEL_PATH):
    global _t5, _t5_tok
    if _t5 is None or _t5_tok is None:
        if not os.path.isdir(model_path):
            raise FileNotFoundError(f"T5 model directory not found: {model_path}")
        # ensure local-only to avoid HF download
        _t5_tok = T5TokenizerFast.from_pretrained(model_path, local_files_only=True)
        _t5 = T5ForConditionalGeneration.from_pretrained(model_path, local_files_only=True)
        _t5.to(device)
        _t5.eval()
    return _t5, _t5_tok

def _clean_t5_output(raw: str) -> str:
    """Small sanitizer for T5 outputs: trim, remove duplicates/artifacts."""
    if not raw:
        return ""
    s = raw.strip()
    # remove repeated short fragments
    s = re.sub(r"(\b\w+\b)(?:\s+\1\b){3,}", r"\1", s)
    # collapse many newlines
    s = re.sub(r"\n\s*\n+", "\n\n", s)
    # strip leading/trailing punctuation fragments
    s = s.strip(" \n\r\t-–—,:;")
    return s

def _t5_quality_check(summary: str) -> bool:
    """
    Return True if summary is acceptable.
    Heuristics: length, meaningful tokens, not too short, not an artifact.
    """
    if not summary or len(summary) < 30:
        return False
    # checks for obvious garbage
    low_content_tokens = ["and solutions", "lorem", "asdf", "no answer", "n/a"]
    summary_low = summary.lower()
    for tok in low_content_tokens:
        if tok in summary_low:
            return False
    # check variety of words
    words = [w for w in re.findall(r"\w+", summary_low) if len(w) > 2]
    if len(words) < 6:
        return False
    return True

def t5_summarize(text: str, max_length: int = 120) -> str:
    """
    Use your fine-tuned T5 for summarization / notes.
    Improved prompt + cleaning + safety checks.
    """
    model, tok = load_t5()
    # keep prompt short and instructive for T5
    # Ask T5 to provide a concise summary in 3 bullets (if possible)
    prompt = f"summarize the text below in a concise way (heading + 3 short bullets, ~50-120 words):\n\n{text}"
    inputs = tok(prompt, return_tensors="pt", truncation=True, padding="longest").to(device)
    with torch.no_grad():
        out = model.generate(**inputs, max_length=max_length, num_beams=4, early_stopping=True)
    raw = tok.decode(out[0], skip_special_tokens=True)
    cleaned = _clean_t5_output(raw)
    return cleaned

def t5_make_notes(text: str, max_length: int = 200, style: str = "bullet") -> str:
    model, tok = load_t5()
    if style == "bullet":
        prompt = f"Create concise bullet notes (short bullets, numbered or bullets):\n\n{text}"
    else:
        prompt = f"Create an outline with headings and subpoints:\n\n{text}"
    inputs = tok(prompt, return_tensors="pt", truncation=True, padding="longest").to(device)
    with torch.no_grad():
        out = model.generate(**inputs, max_length=max_length, num_beams=4, early_stopping=True)
    raw = tok.decode(out[0], skip_special_tokens=True)
    return _clean_t5_output(raw)

# --------------------------- Qwen (local-only) --------------------------
_qwen = None
_qwen_tok = None

def load_qwen(local_path: str = LOCAL_QWEN_PATH):
    """
    Load Qwen from local snapshot. Try low-memory (8-bit) if available,
    then device_map auto, else CPU fallback. Always local_files_only to avoid HF network.
    """
    global _qwen, _qwen_tok
    if _qwen is not None and _qwen_tok is not None:
        return _qwen, _qwen_tok

    if not os.path.isdir(local_path):
        raise FileNotFoundError(f"Local Qwen path not found: {local_path}")

    # tokenizer local-only
    _qwen_tok = AutoTokenizer.from_pretrained(local_path, use_fast=True, local_files_only=True)

    # try bitsandbytes 8-bit (if installed)
    try:
        import bitsandbytes  # noqa: F401
        _qwen = AutoModelForCausalLM.from_pretrained(
            local_path,
            device_map="auto",
            load_in_8bit=True,
            local_files_only=True,
            trust_remote_code=False
        )
        _qwen.eval()
        return _qwen, _qwen_tok
    except Exception:
        pass

    # try device_map auto (for CUDA)
    try:
        _qwen = AutoModelForCausalLM.from_pretrained(
            local_path,
            device_map="auto" if torch.cuda.is_available() else None,
            torch_dtype=torch.float16 if torch.cuda.is_available() else torch.float32,
            local_files_only=True,
            trust_remote_code=False
        )
        _qwen.eval()
        if torch.cuda.is_available():
            _qwen.to("cuda")
        return _qwen, _qwen_tok
    except Exception:
        # final CPU fallback
        _qwen = AutoModelForCausalLM.from_pretrained(
            local_path,
            local_files_only=True,
            torch_dtype=torch.float32,
            trust_remote_code=False
        )
        _qwen.eval()
        _qwen.to("cpu")
        return _qwen, _qwen_tok

def qwen_generate(
    question: str,
    contexts: List[str],
    brief: bool = True,
    max_new_tokens: int = 200,
    temperature: float = 0.25
) -> str:

    model, tok = load_qwen()
    used_ctx = contexts[:3] if contexts else []

    # Build compact context
    ctx_join = "\n\n".join([f"Context {i+1}:\n{clean_text(c)}" for i, c in enumerate(used_ctx)])

    # Detect summary/notes mode
    is_summary_mode = ("summarize" in question.lower()) or ("summary" in question.lower()) or ("notes" in question.lower())

    # Instruction block (not shown to user)
    if brief and not is_summary_mode:
        instruction = (
            "You are an expert study assistant. Answer BRIEFLY using a short heading and 3 bullets. "
            "Use ONLY the provided context. If the answer is not present in context, reply: 'Not found in document.'"
        )
        max_tokens = min(200, max_new_tokens)
    else:
        instruction = (
            "You are an expert study assistant. Provide a clear structured response using only the given text."
        )
        max_tokens = max_new_tokens

    # Use "Answer:" label only for QA mode; summary mode doesn't use "Answer:" to avoid being cut incorrectly
    answer_label = "" if is_summary_mode else "Answer:"

    prompt = f"{instruction}\n\n{ctx_join}\n\nQuestion: {question}\n{answer_label}"

    device = next(model.parameters()).device
    inputs = tok(prompt, return_tensors="pt", truncation=True, padding=True).to(device)

    gen_kwargs = dict(
        max_new_tokens=max_tokens,
        temperature=temperature,
        do_sample=False,
        top_p=0.95,
        repetition_penalty=1.08,
        no_repeat_ngram_size=3,
        eos_token_id=tok.eos_token_id,
        pad_token_id=tok.pad_token_id,
    )

    with torch.no_grad():
        out = model.generate(**inputs, **gen_kwargs)

    raw = tok.decode(out[0], skip_special_tokens=True)

    # -------------------------------------------------------
    #             CLEAN OUTPUT — SHOW ONLY ANSWER
    # -------------------------------------------------------
    cleaned = raw

    if not is_summary_mode:
        # Normal QA mode: extract after "Answer:" if present, otherwise present raw cleaned
        if "Answer:" in cleaned:
            cleaned = cleaned.split("Answer:", 1)[1].strip()
        elif "ANSWER:" in cleaned:
            cleaned = cleaned.split("ANSWER:", 1)[1].strip()
        cleaned = re.sub(r"Context\s*\d+:\s*.*?(?=Answer:|$)", "", cleaned, flags=re.DOTALL)
        cleaned = cleaned.replace("Question:", "").replace("Answer:", "").strip()
        cleaned = re.sub(r"\n\s*\n\s*\n+", "\n\n", cleaned)
        if len(cleaned) < 3:
            cleaned = raw.strip()
        return cleaned
    else:
        # Summary/notes mode: remove prompt echoes gently and return cleaned text
        cleaned = re.sub(r"Question:.*", "", cleaned, flags=re.DOTALL)
        cleaned = re.sub(r"Context\s*\d+:", "", cleaned)
        cleaned = cleaned.strip()
        cleaned = re.sub(r"\n\s*\n+", "\n\n", cleaned)
        # remove trailing small junk words
        cleaned = re.sub(r"(?i)\b(and|solutions?|notes?)\s*$", "", cleaned).strip()
        if len(cleaned) < 3:
            cleaned = raw.strip()
        return cleaned

# ----------------------- High-level pipeline ----------------------------
def ask_with_rag(question: str, index, metas: List[Dict], top_k: int = 5, brief: bool = True) -> str:
    ctxs = retrieve_contexts(index, metas, question, top_k=top_k, min_score=0.05)
    if not ctxs:
        return "No relevant information found in the uploaded document."
    return qwen_generate(question, ctxs, brief=brief, max_new_tokens=180, temperature=0.25)

def ask_without_rag(question: str, prefer_brief: bool = True) -> str:
    """
    Non-document questions: use Qwen for best QA quality.
    """
    try:
        return qwen_generate(question, contexts=[], brief=prefer_brief, max_new_tokens=180, temperature=0.25)
    except Exception as e:
        return f"Error generating answer: {str(e)}"

# ---------------- Summarize & Notes exposed (T5-backed with Hybrid) ------
def summarize_text(text: str, max_length: int = 120, hybrid_auto: bool = True) -> str:
    """
    Hybrid summarization:
      1) Run T5 (project requirement)
      2) Validate T5 output (heuristic)
      3a) If acceptable -> optionally refine with Qwen (hybrid_auto param controls auto-refinement)
      3b) If NOT acceptable -> ask Qwen to produce/refine a high-quality summary from the text
    hybrid_auto: If True, Qwen refinement runs only when T5 is poor OR to polish T5 automatically (auto polish)
                 If False, return T5 summary if T5 passes checks.
    """
    # 1) T5 attempt
    try:
        t5_out = t5_summarize(text, max_length=max_length)
    except Exception as e:
        t5_out = ""

    t5_out = _clean_t5_output(t5_out)

    # 2) Check T5 quality
    t5_ok = _t5_quality_check(t5_out)

    # 3) Hybrid logic (Option C = AUTO)
    try:
        if t5_ok:
            # T5 produced a decent summary
            if hybrid_auto:
                # Ask Qwen to polish the T5 output (rewrite into heading + 3 bullets)
                refine_prompt = f"Rewrite and polish this summary into a concise structured summary (heading + 3 bullets, clear, brief):\n\n{t5_out}"
                qwen_refined = qwen_generate(refine_prompt, [text], brief=True, max_new_tokens=max_length, temperature=0.2)
                qwen_refined = qwen_refined.strip()
                # If Qwen refinement produced something reasonable, return it; otherwise return T5
                if qwen_refined and len(qwen_refined) > 20:
                    return qwen_refined
                else:
                    return t5_out
            else:
                # hybrid_auto disabled -> return T5 output
                return t5_out
        else:
            # T5 failed quality check -> use Qwen to create summary directly
            qwen_prompt = "Summarize the following text into a concise heading and 3 short bullets. Provide a very clear brief summary."
            qwen_out = qwen_generate(qwen_prompt + "\n\n" + text, [text], brief=True, max_new_tokens=max_length, temperature=0.2)
            qwen_out = qwen_out.strip()
            # final fallback: if qwen_out is empty, try to return t5_out (even if poor)
            if not qwen_out or len(qwen_out) < 10:
                return t5_out or "No summary could be generated."
            return qwen_out
    except Exception as e:
        # in case Qwen fails, fallback to T5 (even if poor) with some cleanup
        return t5_out or f"Error generating summary: {str(e)}"

def generate_notes(text: str, style: str = "bullet", max_length: int = 200, hybrid_auto: bool = True) -> str:
    """
    Notes pipeline:
      - Try T5 notes first (project requirement)
      - If T5 notes are poor or empty, use Qwen to generate notes from the source
      - If T5 is okay and hybrid_auto True, Qwen can polish the notes.
    """
    try:
        t5_out = t5_make_notes(text, max_length=max_length, style=style)
    except Exception:
        t5_out = ""
    t5_out = _clean_t5_output(t5_out)
    t5_ok = _t5_quality_check(t5_out)

    try:
        if t5_ok:
            if hybrid_auto:
                # polish with Qwen (short bullets)
                prompt = f"Polish and format these study notes into concise bullets (max 12 bullets):\n\n{t5_out}"
                polished = qwen_generate(prompt, [text], brief=True, max_new_tokens=max_length, temperature=0.2)
                polished = polished.strip()
                return polished or t5_out
            return t5_out
        else:
            # T5 poor -> generate with Qwen directly
            prompt = f"Create concise study notes from the text below in bullets (max 12 bullets):\n\n{text}"
            qwen_out = qwen_generate(prompt, [text], brief=True, max_new_tokens=max_length, temperature=0.2)
            qwen_out = qwen_out.strip()
            return qwen_out or t5_out or "No notes could be generated."
    except Exception as e:
        return t5_out or f"Error generating notes: {str(e)}"

# ---------------- Document ingestion (index + metas) ---------------------
def ingest_document(path: str, index_path: str, meta_path: str, chunk_size: int = 400, overlap: int = 64):
    p = Path(path)
    ext = p.suffix.lower()
    if ext == ".pdf":
        text = extract_text_from_pdf(path)
    elif ext == ".docx":
        text = extract_text_from_docx(path)
    elif ext == ".pptx":
        text = extract_text_from_pptx(path)
    elif ext == ".txt":
        text = extract_text_from_txt(path)
    elif ext in [".png", ".jpg", ".jpeg", ".tiff", ".bmp"]:
        text = extract_text_from_image(path)
    else:
        raise ValueError("Unsupported file type: " + ext)

    # Clean & de-duplicate repeated paragraph blocks (helpful for broken PDFs)
    text = clean_text(text)
    # Simple de-duplication heuristic: remove repeated consecutive paragraphs
    paras = [p.strip() for p in re.split(r'\n{1,}', text) if p.strip()]
    compact = []
    prev = None
    for p in paras:
        if p == prev:
            continue
        compact.append(p)
        prev = p
    text = "\n\n".join(compact)

    chunks = chunk_text(text, chunk_size=chunk_size, overlap=overlap)

    metas = []
    doc_id = str(uuid.uuid4())
    for i, c in enumerate(chunks):
        metas.append({"doc_id": doc_id, "chunk_id": f"{doc_id}_{i}", "text": c, "source": str(path)})

    # build index and write metas
    idx = build_index([m["text"] for m in metas], index_path)
    os.makedirs(os.path.dirname(meta_path), exist_ok=True)
    with open(meta_path, "w", encoding="utf-8") as f:
        json.dump(metas, f, indent=2)
    return idx, metas

# ---------------- Unified GUI entrypoint -------------------------------
def full_ai_pipeline(question: str,
                     use_rag: bool = False,
                     index = None,
                     metas: Optional[List[Dict]] = None,
                     top_k: int = 5,
                     brief: bool = True) -> str:
    """
    If use_rag True and index/metas provided -> use RAG + Qwen
    Otherwise use Qwen for best QA quality (T5 reserved for summaries/notes)
    """
    try:
        if use_rag and index is not None and metas is not None:
            return ask_with_rag(question, index, metas, top_k=top_k, brief=brief)
        else:
            return ask_without_rag(question, prefer_brief=brief)
    except Exception as e:
        return f"Error in backend pipeline: {str(e)}"

# ------------------------- Health check --------------------------------
def health_check() -> Dict[str, str]:
    status = {"device": device}
    # t5 status
    try:
        status["t5"] = "available" if os.path.isdir(T5_MODEL_PATH) else "missing"
    except Exception as e:
        status["t5"] = f"error: {str(e)}"
    # qwen tokenizer check (local only)
    try:
        _ = AutoTokenizer.from_pretrained(LOCAL_QWEN_PATH, use_fast=True, local_files_only=True)
        status["qwen_tokenizer"] = "available"
    except Exception as e:
        status["qwen_tokenizer"] = f"not available: {str(e)}"
    status["embeddings"] = "available" if SentenceTransformer is not None else "missing"
    status["faiss"] = "available" if faiss is not None else "missing"
    return status

# End of file
